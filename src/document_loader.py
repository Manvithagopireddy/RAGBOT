"""
document_loader.py — Multi-format document loader supporting PDF, DOCX, TXT, CSV, PPTX.
"""
import os
import csv
import io
from pathlib import Path
from typing import List
# pyright: ignore [missing-import]
from langchain_core.documents import Document
# pyright: ignore [missing-import]
from langchain_community.document_loaders import PyPDFLoader
from src.logger import get_logger

logger = get_logger(__name__)


def load_pdf_file(file_path: Path) -> List[Document]:
    """Loads a PDF file and returns its content as a list of Documents."""
    logger.info(f"Loading PDF: {file_path.name}")
    try:
        loader = PyPDFLoader(str(file_path))
        documents = loader.load()
        for doc in documents:
            doc.metadata["source"] = file_path.name
        logger.info(f"Loaded {len(documents)} pages from {file_path.name}")
        return documents
    except Exception as e:
        logger.error(f"Failed to load PDF {file_path.name}: {e}")
        return []


def load_txt_file(file_path: Path) -> List[Document]:
    """Loads a plain text file."""
    try:
        text = file_path.read_text(encoding="utf-8", errors="ignore")
        return [Document(page_content=text, metadata={"source": file_path.name, "page": 1})]
    except Exception as e:
        logger.error(f"Failed to load TXT {file_path.name}: {e}")
        return []


def load_csv_file(file_path: Path) -> List[Document]:
    """Loads a CSV file, converting rows to text."""
    try:
        rows = []
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.DictReader(f)
            for i, row in enumerate(reader):
                row_text = " | ".join(f"{k}: {v}" for k, v in row.items())
                rows.append(Document(
                    page_content=row_text,
                    metadata={"source": file_path.name, "page": i + 1}
                ))
        return rows
    except Exception as e:
        logger.error(f"Failed to load CSV {file_path.name}: {e}")
        return []


def load_docx_file(file_path: Path) -> List[Document]:
    """Loads a Microsoft Word (.docx) file."""
    try:
        import docx  # python-docx
        doc = docx.Document(str(file_path))
        full_text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
        return [Document(page_content=full_text, metadata={"source": file_path.name, "page": 1})]
    except ImportError:
        logger.warning("python-docx not installed. Install with: pip install python-docx")
        return []
    except Exception as e:
        logger.error(f"Failed to load DOCX {file_path.name}: {e}")
        return []


def load_pptx_file(file_path: Path) -> List[Document]:
    """Loads a Microsoft PowerPoint (.pptx) file."""
    try:
        from pptx import Presentation  # python-pptx
        prs = Presentation(str(file_path))
        docs = []
        for i, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
            if text_parts:
                docs.append(Document(
                    page_content="\n".join(text_parts),
                    metadata={"source": file_path.name, "page": i + 1}
                ))
        return docs
    except ImportError:
        logger.warning("python-pptx not installed. Install with: pip install python-pptx")
        return []
    except Exception as e:
        logger.error(f"Failed to load PPTX {file_path.name}: {e}")
        return []


def load_file(file_path: Path) -> List[Document]:
    """Dispatches loading to the correct loader based on file extension."""
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return load_pdf_file(file_path)
    elif ext == ".txt":
        return load_txt_file(file_path)
    elif ext == ".csv":
        return load_csv_file(file_path)
    elif ext == ".docx":
        return load_docx_file(file_path)
    elif ext == ".pptx":
        return load_pptx_file(file_path)
    else:
        logger.warning(f"Unsupported file type: {ext}")
        return []


def load_all_documents(directory_path: Path) -> List[Document]:
    """Scans a directory for all supported files and loads them."""
    logger.info(f"Scanning directory: {directory_path}")
    if not directory_path.exists():
        directory_path.mkdir(parents=True, exist_ok=True)
        return []

    supported_extensions = {".pdf", ".txt", ".csv", ".docx", ".pptx"}
    files = [
        f for f in directory_path.iterdir()
        if f.is_file() and f.suffix.lower() in supported_extensions
    ]

    if not files:
        logger.warning(f"No supported files found in {directory_path}")
        return []

    logger.info(f"Found {len(files)} files to load.")
    all_documents = []
    for f in files:
        all_documents.extend(load_file(f))

    logger.info(f"Total chunks loaded: {len(all_documents)}")
    return all_documents
